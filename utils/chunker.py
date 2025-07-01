import os
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from typing import List
import re

# Supported file extensions
alLOWED_EXT = {".pdf", ".txt", ".pptx", ".doc", ".docx"}

def extract_text(path: str) -> str:
    """Extract text from PDF, PPTX, DOCX/DOC, or TXT files with validation"""
    if not os.path.exists(path):
        raise FileNotFoundError(f"File not found: {path}")
    ext = os.path.splitext(path)[1].lower()
    if ext not in alLOWED_EXT:
        raise ValueError("Unsupported file type. Only PDF, PPTX, DOCX, and TXT are supported.")

    text_segments: List[str] = []
    try:
        if ext == ".pdf":
            reader = PdfReader(path)
            for page in reader.pages:
                text_segments.append(page.extract_text() or "")
        elif ext == ".pptx":
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_segments.append(shape.text)
        elif ext in {".doc", ".docx"}:
            doc = Document(path)
            for para in doc.paragraphs:
                text_segments.append(para.text)
        else:  # .txt
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
    except Exception as e:
        raise ValueError(f"Extraction failed for {ext.upper()}: {str(e)}")

    # Join with line breaks
    return "\n".join(text_segments)


def chunk_text(text: str, size: int = 384, overlap: int = 64) -> List[str]:
    """Chunk text into overlapping segments with sentence-level splits"""
    # Split into sentences
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.) (?<=\.|\?|\!)\s', text)
    chunks: List[str] = []
    current: List[str] = []
    current_len = 0

    for sent in sentences:
        sent_words = sent.split()
        sent_len = len(sent_words)
        if current_len + sent_len > size and current:
            chunks.append(" ".join(current))
            # Retain overlap words
            if overlap > 0:
                tail = " ".join(current).split()[-overlap:]
                current = tail
                current_len = len(tail)
            else:
                current = []
                current_len = 0
        current.append(sent)
        current_len += sent_len

    if current:
        chunks.append(" ".join(current))
    return chunks
